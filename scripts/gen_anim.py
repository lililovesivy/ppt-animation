#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT动画生成器 v5 - 智能元素分类 + 多样化动画
基于原生XML修改方案，确保PowerPoint兼容性
"""

import os
import zipfile
import shutil
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# ============================================================
# 配置：元素类型与动画映射规则
# ============================================================

# 飞入方向映射 (presetSubtype)
FLY_DIRECTIONS = {
    'left': 4,       # 自左
    'right': 2,      # 自右
    'bottom': 3,     # 自下
    'top': 1,        # 自上
    'left-top': 6,   # 左上
    'right-top': 7,  # 右上
    'right-bottom': 8,  # 右下
    'left-bottom': 9,   # 左下
}

# 擦除方向映射 (presetSubtype)
WIPE_DIRECTIONS = {
    'left': 4,       # 自左
    'right': 2,      # 自右
    'bottom': 3,     # 自下
    'top': 1,        # 自上
}

# 元素类型动画配置
ELEMENT_ANIMATION_CONFIG = {
    'bg_image': {
        'animation': 'none',  # 背景图无动画，固定显示
        'direction': 'none',
        'trigger': 'none',
        'description': '背景图（固定显示，无动画）',
    },
    'title_text': {
        'animation': 'wipe',
        'direction': 'left',
        'trigger': 'on_click',
        'build_text': 'all_at_once',
        'description': '标题文字',
    },
    'body_text': {
        'animation': 'wipe',
        'direction': 'left',
        'trigger': 'on_click',
        'build_text': 'by_paragraph',
        'description': '正文文字',
    },
    'emphasis_text': {
        'animation': 'fly_in',
        'direction': 'bottom',
        'trigger': 'on_click',
        'build_text': 'all_at_once',
        'description': '重点文字',
    },
    'deco_shape': {
        'animation': 'fly_in',
        'direction': 'auto',  # 根据位置自动决定
        'trigger': 'with_previous',
        'description': '装饰形状',
    },
    'highlight_shape': {
        'animation': 'wipe',
        'direction': 'left',
        'trigger': 'before_previous',
        'description': '高亮框',
    },
    'content_image': {
        'animation': 'fly_in',
        'direction': 'auto',
        'trigger': 'on_click',
        'description': '内容图片',
    },
    'icon_image': {
        'animation': 'fly_in',
        'direction': 'auto',
        'trigger': 'with_previous',
        'description': '小图标',
    },
    'line_shape': {
        'animation': 'wipe',
        'direction': 'left',
        'trigger': 'on_click',
        'description': '线条',
    },
    'unknown': {
        'animation': 'fly_in',
        'direction': 'auto',
        'trigger': 'on_click',
        'description': '未知元素',
    },
}


# ============================================================
# 元素分类函数
# ============================================================

def classify_shape(shape, slide_width, slide_height):
    """
    智能分类元素类型
    返回: (类型标签, 方向建议)
    """
    shape_type = shape.shape_type
    has_text = shape.has_text_frame and shape.text_frame.text.strip()
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    
    # 计算相对位置 (0-1)
    rel_x = left / slide_width if slide_width > 0 else 0.5
    rel_y = top / slide_height if slide_height > 0 else 0.5
    rel_w = width / slide_width if slide_width > 0 else 0
    rel_h = height / slide_height if slide_height > 0 else 0
    
    # 1. 判断是否为背景图
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        if rel_w > 0.9 and rel_h > 0.9 and abs(rel_x) < 0.05 and abs(rel_y) < 0.05:
            direction = get_direction_by_position(rel_x, rel_y, 'image')
            return 'bg_image', direction
        
        # 小图标
        if rel_w < 0.15 and rel_h < 0.15:
            direction = get_direction_by_position(rel_x, rel_y, 'image')
            return 'icon_image', direction
        
        # 普通内容图片
        direction = get_direction_by_position(rel_x, rel_y, 'image')
        return 'content_image', direction
    
    # 2. 判断是否为线条
    if shape_type == MSO_SHAPE_TYPE.LINE:
        if width > height * 3:  # 水平线
            return 'line_shape', 'left'
        else:  # 垂直线
            return 'line_shape', 'top'
    
    # 细长形状也视为线条（但必须没有文字，且绝对尺寸很小）
    if not has_text:
        is_very_thin = (width > height * 8 and height < 200000) or (height > width * 8 and width < 200000)
        if is_very_thin:
            if width > height:
                return 'line_shape', 'left'
            else:
                return 'line_shape', 'top'
    
    # 3. 文本相关
    if has_text:
        text = shape.text_frame.text
        
        # 检查是否为标题（文字少、字号大）
        if len(text) < 30 and len(text) > 0:
            # 尝试获取字号
            try:
                if shape.text_frame.paragraphs:
                    first_run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                    if first_run and first_run.font.size:
                        font_size = first_run.font.size.pt
                        if font_size >= 24:
                            return 'title_text', 'left'
            except:
                pass
            
            # 位于页面上方的短文本视为标题
            if rel_y < 0.3:
                return 'title_text', 'left'
        
        # 检查是否为重点文字（红色）
        try:
            if shape.text_frame.paragraphs:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.type is not None:
                            # 红色系判断
                            rgb = str(run.font.color.rgb) if hasattr(run.font.color, 'rgb') else ''
                            if rgb and rgb.startswith('FF') or 'red' in str(run.font.color.type).lower():
                                return 'emphasis_text', 'bottom'
        except:
            pass
        
        # 正文文字
        return 'body_text', 'left'
    
    # 4. 纯形状
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        # 高亮框（较宽较矮，位于文字区域下方）
        if rel_w > 0.3 and rel_h < 0.15:
            return 'highlight_shape', 'left'
        
        # 小装饰形状
        if rel_w < 0.1 and rel_h < 0.1:
            direction = get_direction_by_position(rel_x, rel_y, 'shape')
            return 'deco_shape', direction
        
        # 普通形状
        direction = get_direction_by_position(rel_x, rel_y, 'shape')
        return 'deco_shape', direction
    
    # 5. 文本框
    if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        if has_text:
            text = shape.text_frame.text
            if len(text) < 30 and rel_y < 0.3:
                return 'title_text', 'left'
            return 'body_text', 'left'
    
    # 默认
    direction = get_direction_by_position(rel_x, rel_y, 'default')
    return 'unknown', direction


def get_direction_by_position(rel_x, rel_y, element_type='default'):
    """
    根据元素在页面中的位置，智能决定动画方向
    让元素从靠近的边缘飞入，更自然
    """
    # 计算距离各边的距离
    dist_left = rel_x
    dist_right = 1 - rel_x
    dist_top = rel_y
    dist_bottom = 1 - rel_y
    
    # 找最近的边
    distances = {
        'left': dist_left,
        'right': dist_right,
        'top': dist_top,
        'bottom': dist_bottom,
    }
    min_dir = min(distances, key=distances.get)
    min_dist = distances[min_dir]
    
    # 对于图片，使用对角线方向增加动感
    if element_type in ['image', 'icon_image', 'content_image', 'deco_shape']:
        # 判断是否靠近角落
        if min_dist < 0.2:
            # 找第二近的边
            sorted_dirs = sorted(distances, key=distances.get)
            second_dir = sorted_dirs[1]
            if distances[second_dir] < 0.3:
                # 组合成对角线方向
                diagonal_map = {
                    ('left', 'top'): 'left-top',
                    ('top', 'left'): 'left-top',
                    ('right', 'top'): 'right-top',
                    ('top', 'right'): 'right-top',
                    ('left', 'bottom'): 'left-bottom',
                    ('bottom', 'left'): 'left-bottom',
                    ('right', 'bottom'): 'right-bottom',
                    ('bottom', 'right'): 'right-bottom',
                }
                diag = diagonal_map.get((min_dir, second_dir))
                if diag:
                    return diag
    
    return min_dir


# ============================================================
# XML处理函数
# ============================================================

def find_matching_par_end(content, start_pos):
    """从start_pos（<p:par>之后）找匹配的</p:par>，返回结束位置"""
    depth = 1
    pos = start_pos
    while depth > 0 and pos < len(content):
        next_open = content.find('<p:par>', pos)
        next_close = content.find('</p:par>', pos)
        if next_close == -1:
            break
        if next_open != -1 and next_open < next_close:
            depth += 1
            pos = next_open + len('<p:par>')
        else:
            depth -= 1
            pos = next_close + len('</p:par>')
    return pos if depth == 0 else -1


def extract_anim_templates(base_timing):
    """
    从基准timing中提取动画模板
    返回: {
        'fly_in': 飞入模板字符串,
        'wipe': 擦除模板(无grpId),
        'wipe_text': 擦除模板(有grpId),
        'prefix': 动画前的内容,
        'suffix': 动画后的内容,
    }
    """
    # 找所有 delay="indefinite" 的位置
    indef_positions = [m.start() for m in re.finditer(r'delay="indefinite"', base_timing)]
    
    anim_blocks = []
    
    for pos in indef_positions:
        # 向前找这个cTn所在的par
        ctn_start = base_timing.rfind('<p:cTn', 0, pos)
        if ctn_start == -1:
            continue
        
        # 检查有没有nodeType
        ctn_tag_end = base_timing.find('>', ctn_start)
        ctn_header = base_timing[ctn_start:ctn_tag_end+1]
        if 'nodeType=' in ctn_header and 'tmRoot' in ctn_header:
            continue  # 跳过tmRoot
        
        # 向前找最近的<p:par>（外层par）
        par_start = base_timing.rfind('<p:par>', 0, ctn_start)
        if par_start == -1:
            continue
        
        # 找匹配的</p:par>
        par_end = find_matching_par_end(base_timing, par_start + len('<p:par>'))
        if par_end == -1:
            continue
        
        # 确认包含clickEffect
        block = base_timing[par_start:par_end]
        if 'nodeType="clickEffect"' not in block:
            continue
        
        anim_blocks.append(block)
    
    if len(anim_blocks) < 2:
        raise Exception(f"提取动画模板失败，只找到 {len(anim_blocks)} 个动画块")
    
    # 分类模板
    templates = {
        'fly_in': None,
        'wipe': None,
        'wipe_text': None,
    }
    
    for blk in anim_blocks:
        has_grp = 'grpId=' in blk
        preset_match = re.search(r'presetID="(\d+)"', blk)
        preset_id = preset_match.group(1) if preset_match else None
        
        if preset_id == '2' and not has_grp:
            templates['fly_in'] = blk
        elif preset_id == '22' and has_grp:
            templates['wipe_text'] = blk
        elif preset_id == '22' and not has_grp:
            templates['wipe'] = blk
    
    # 确保至少有两个模板可用
    if not templates['fly_in']:
        templates['fly_in'] = anim_blocks[0]
    if not templates['wipe']:
        templates['wipe'] = anim_blocks[1] if len(anim_blocks) > 1 else anim_blocks[0]
    if not templates['wipe_text']:
        templates['wipe_text'] = templates['wipe']
    
    # 找prefix和suffix
    first_block_start = base_timing.find(anim_blocks[0])
    last_block_end = base_timing.find(anim_blocks[-1]) + len(anim_blocks[-1])
    
    templates['prefix'] = base_timing[:first_block_start]
    templates['suffix'] = base_timing[last_block_end:]
    
    return templates


def modify_anim_block(template, shape_id, preset_id, preset_subtype, grp_id=None):
    """修改动画块的属性"""
    result = template
    
    # 替换spid
    result = re.sub(r'spid="[^"]*"', f'spid="{shape_id}"', result)
    
    # 替换presetID
    result = re.sub(r'presetID="[^"]*"', f'presetID="{preset_id}"', result)
    
    # 替换presetSubtype
    result = re.sub(r'presetSubtype="[^"]*"', f'presetSubtype="{preset_subtype}"', result)
    
    # 替换animEffect的filter属性
    if '<p:animEffect' in result:
        # 从filter映射表找对应的值
        filter_map = {
            'left': 'wipe(left)',
            'right': 'wipe(right)',
            'top': 'wipe(top)',
            'bottom': 'wipe(bottom)',
        }
        # 找当前filter类型
        filter_match = re.search(r'filter="wipe\(([^)]+)\)"', result)
        if filter_match:
            current_dir = filter_match.group(1)
            # 根据preset_subtype反查方向
            dir_map = {v: k for k, v in WIPE_DIRECTIONS.items()}
            target_dir = dir_map.get(preset_subtype, current_dir)
            result = re.sub(
                r'filter="wipe\([^)]+\)"',
                f'filter="wipe({target_dir})"',
                result
            )
    
    # 处理grpId
    if grp_id is not None:
        # 移除已有的grpId
        result = re.sub(r'\s+grpId="[^"]*"', '', result)
        # 在nodeType="clickEffect"的cTn上添加grpId
        result = re.sub(
            r'(nodeType="clickEffect")',
            f'grpId="{grp_id}" \\1',
            result
        )
    else:
        # 移除grpId
        result = re.sub(r'\s+grpId="[^"]*"', '', result)
    
    return result


def generate_bldLst(text_shapes_info):
    """生成构建列表，控制文本动画的构建方式"""
    entries = []
    for shape_id, build_type in text_shapes_info:
        if build_type == 'by_paragraph':
            # 按段落构建 - 每个段落一个动画
            # 0级是整体，1级是第一段落，以此类推
            entries.append(f'<p:bldP spid="{shape_id}" grpId="0" animBg="1"/>')
            entries.append(f'<p:bldP spid="{shape_id}" grpId="1" animBg="1"/>')
        else:
            # 整体出现
            entries.append(f'<p:bldP spid="{shape_id}" grpId="0" animBg="1"/>')
    
    return '<p:bldLst>' + ''.join(entries) + '</p:bldLst>'


# ============================================================
# 动画生成主函数
# ============================================================

def generate_animations_for_slide(shapes_info, templates, slide_idx):
    """
    为单页幻灯片生成动画
    shapes_info: [(shape_id, category, direction, has_text, build_type), ...]
    """
    new_blocks = []
    text_shapes_for_bld = []  # [(shape_id, build_type), ...]
    
    for i, (shape_id, category, direction, has_text, build_type) in enumerate(shapes_info):
        config = ELEMENT_ANIMATION_CONFIG.get(category, ELEMENT_ANIMATION_CONFIG['unknown'])
        
        # 跳过无动画的元素（如背景图）
        anim_type = config['animation']
        if anim_type == 'none':
            continue
        
        # 决定方向
        if direction == 'auto' or config['direction'] != 'auto':
            if config['direction'] != 'auto':
                direction = config['direction']
        
        # 选择模板
        if anim_type in ['fly_in', 'fade']:
            template = templates['fly_in']
            preset_id = 2
            subtype = FLY_DIRECTIONS.get(direction, 4)
            grp_id = None
        elif anim_type == 'wipe':
            if has_text and build_type == 'by_paragraph':
                template = templates['wipe_text']
                grp_id = 0
                text_shapes_for_bld.append((shape_id, build_type))
            elif has_text:
                template = templates['wipe_text']
                grp_id = 0
                text_shapes_for_bld.append((shape_id, build_type))
            else:
                template = templates['wipe']
                grp_id = None
            
            preset_id = 22
            subtype = WIPE_DIRECTIONS.get(direction, 4)
        else:
            # 默认用飞入
            template = templates['fly_in']
            preset_id = 2
            subtype = FLY_DIRECTIONS.get(direction, 4)
            grp_id = None
        
        # 修改动画块
        new_block = modify_anim_block(
            template,
            shape_id,
            preset_id,
            subtype,
            grp_id
        )
        new_blocks.append(new_block)
    
    # 构建新的timing
    new_timing_body = templates['prefix'] + ''.join(new_blocks) + templates['suffix']
    
    # 更新bldLst
    if text_shapes_for_bld:
        new_bldLst = generate_bldLst(text_shapes_for_bld)
        if '<p:bldLst/>' in new_timing_body:
            new_timing_body = new_timing_body.replace('<p:bldLst/>', new_bldLst)
        elif '<p:bldLst>' in new_timing_body:
            new_timing_body = re.sub(r'<p:bldLst>.*?</p:bldLst>', new_bldLst, new_timing_body, flags=re.DOTALL)
        else:
            new_timing_body = new_timing_body.replace('</p:timing>', new_bldLst + '</p:timing>')
    
    return new_timing_body


def sort_shapes_by_visual_order(shapes_info, slide_width, slide_height):
    """
    按视觉层次排序元素出场顺序
    规则：从上到下，从左到右，背景先，标题先，图片后
    """
    def sort_key(item):
        shape_id, category, direction, has_text, build_type, left, top, width, height = item
        
        # 优先级分数（越低越先出现）
        priority = 0
        
        # 1. 背景最先
        if category == 'bg_image':
            priority += 0
        # 2. 标题其次
        elif category == 'title_text':
            priority += 100
        # 3. 高亮框在文字前
        elif category == 'highlight_shape':
            priority += 200
        # 4. 正文文字
        elif category in ['body_text', 'emphasis_text']:
            priority += 300
        # 5. 图片
        elif category in ['content_image', 'icon_image']:
            priority += 400
        # 6. 装饰形状
        elif category in ['deco_shape', 'line_shape']:
            priority += 500
        else:
            priority += 600
        
        # 按位置排序（从上到下，从左到右）
        priority += top // 10000  # 垂直位置（越靠上越先）
        priority += left // 100000  # 水平位置（越靠左越先）
        
        return priority
    
    return sorted(shapes_info, key=sort_key)


# ============================================================
# 主流程
# ============================================================

def main():
    base_dir = '/app/data/所有对话/主对话'
    original_pptx = os.path.join(base_dir, '用户上传/副本闪婚闪离-13-16-可编辑_1781581419182_0_mp3e.pptx')
    output_pptx = os.path.join(base_dir, '闪婚闪离-智能动画版.pptx')
    temp_dir = os.path.join(base_dir, 'ppt_temp_anim_v5')
    
    # 清理临时目录
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
    os.makedirs(temp_dir)
    
    # 解压原始PPT
    print("解压PPT...")
    with zipfile.ZipFile(original_pptx, 'r') as zf:
        zf.extractall(temp_dir)
    
    # 读取slide2的timing作为模板
    slide2_xml_path = os.path.join(temp_dir, 'ppt/slides/slide2.xml')
    with open(slide2_xml_path, 'r', encoding='utf-8') as f:
        slide2_content = f.read()
    
    timing_match = re.search(r'<p:timing>.*?</p:timing>', slide2_content, re.DOTALL)
    if not timing_match:
        print("错误：找不到timing")
        return
    base_timing = timing_match.group(0)
    
    # 提取动画模板
    print("提取动画模板...")
    templates = extract_anim_templates(base_timing)
    print(f"  飞入模板: {'✓' if templates['fly_in'] else '✗'}")
    print(f"  擦除模板(无grp): {'✓' if templates['wipe'] else '✗'}")
    print(f"  擦除模板(有grp): {'✓' if templates['wipe_text'] else '✗'}")
    
    # 获取总页数和页面尺寸
    prs = Presentation(original_pptx)
    num_slides = len(prs.slides)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    print(f"\n共 {num_slides} 页幻灯片")
    
    # 为每一页添加动画
    for slide_idx in range(num_slides):
        print(f"\n=== 第 {slide_idx+1} 页 ===")
        
        slide = prs.slides[slide_idx]
        shapes_info = []
        
        # 分析每个元素
        for shape in slide.shapes:
            category, direction = classify_shape(shape, slide_width, slide_height)
            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            
            # 决定文本构建方式
            build_type = 'all_at_once'
            if category in ['body_text'] and has_text:
                # 有多段文字的用by_paragraph
                try:
                    if len(shape.text_frame.paragraphs) > 1:
                        build_type = 'by_paragraph'
                except:
                    pass
            
            shapes_info.append((
                shape.shape_id, category, direction, has_text, build_type,
                shape.left, shape.top, shape.width, shape.height
            ))
        
        # 按视觉顺序排序
        shapes_info_sorted = sort_shapes_by_visual_order(shapes_info, slide_width, slide_height)
        
        # 输出分类结果
        print(f"  元素总数: {len(shapes_info_sorted)}")
        for i, (sid, cat, direc, has_txt, build, l, t, w, h) in enumerate(shapes_info_sorted):
            config = ELEMENT_ANIMATION_CONFIG.get(cat, {})
            desc = config.get('description', cat)
            anim_type = config.get('animation', '?')
            print(f"    [{i+1}] id={sid}, {desc}, 动画={anim_type}, 方向={direc}")
        
        # 生成动画
        # 简化参数，只传需要的
        anim_input = [(sid, cat, direc, has_txt, build) 
                     for sid, cat, direc, has_txt, build, l, t, w, h in shapes_info_sorted]
        
        new_timing = generate_animations_for_slide(anim_input, templates, slide_idx)
        
        # 写入幻灯片XML
        slide_xml_path = os.path.join(temp_dir, f'ppt/slides/slide{slide_idx+1}.xml')
        with open(slide_xml_path, 'r', encoding='utf-8') as f:
            slide_content = f.read()
        
        # 替换或添加timing
        if '<p:timing>' in slide_content:
            slide_content = re.sub(r'<p:timing>.*?</p:timing>', new_timing, slide_content, flags=re.DOTALL)
        else:
            slide_content = slide_content.replace('</p:sld>', new_timing + '</p:sld>')
        
        with open(slide_xml_path, 'w', encoding='utf-8') as f:
            f.write(slide_content)
    
    # 重新打包PPT
    print("\n打包PPT...")
    with zipfile.ZipFile(output_pptx, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, temp_dir)
                zf.write(file_path, arcname)
    
    # 清理
    shutil.rmtree(temp_dir)
    
    print(f"\n✓ 完成！输出文件: {output_pptx}")
    print(f"  共 {num_slides} 页，全部添加了智能分类动画")


if __name__ == '__main__':
    main()
